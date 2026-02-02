import cv2
import os
import numpy as np
from pptx import Presentation
from pptx.util import Inches

# ================= 配置区 =================
INTERVAL_SECONDS = 1      # 采样频率 (秒)
SENSITIVITY_FACTOR = 1.3  # 后期突变判定倍数
EARLY_THRESHOLD = 8000    # 前150秒硬阈值
OUTPUT_FILENAME = "全集_精炼汇总版.pptx" # 最终输出的文件名
# ==========================================

def add_slide(prs, frame):
    """
    保存当前帧到PPT，并立即清理临时文件
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    temp_img = "temp_capture.jpg"
    cv2.imwrite(temp_img, frame)
    slide.shapes.add_picture(temp_img, 0, 0, width=Inches(10))
    if os.path.exists(temp_img):
        os.remove(temp_img)

def process_one_video(video_path, prs, video_index, total_videos):
    """
    处理单个视频，将结果追加到 prs 对象中
    """
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print(f"❌ 无法打开视频: {video_path}")
        return 0

    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    frame_interval = int(fps * INTERVAL_SECONDS)
    
    last_gray = None 
    pending_frame = None   
    diff_history = []      
    local_slide_count = 0 # 当前视频截了多少页
    count = 0

    print(f"\n" + "="*65)
    print(f"🎬 [{video_index}/{total_videos}] 正在收割: {video_path}")
    print("="*65)

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret: 
            if pending_frame is not None:
                add_slide(prs, pending_frame)
                local_slide_count += 1
            break
        
        if count % frame_interval == 0:
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            gray_blur = cv2.GaussianBlur(gray, (21, 21), 0)
            
            if last_gray is not None:
                delta = cv2.absdiff(last_gray, gray_blur)
                thresh = cv2.threshold(delta, 25, 255, cv2.THRESH_BINARY)[1]
                current_diff = cv2.countNonZero(thresh)
                
                avg_diff = np.mean(diff_history) if len(diff_history) > 0 else current_diff
                current_ratio = current_diff / avg_diff if avg_diff > 0 else 0
                
                is_triggered = False

                # 判定逻辑
                if count < fps * 150:
                    if current_diff > EARLY_THRESHOLD:
                        is_triggered = True
                else:
                    if len(diff_history) > 3 and current_ratio > SENSITIVITY_FACTOR:
                        is_triggered = True

                if is_triggered:
                    if pending_frame is not None:
                        add_slide(prs, pending_frame)
                        local_slide_count += 1
                        
                        timestamp = f"{count//int(fps)}s"
                        print(f"\r✨ [本视频第 {local_slide_count:02d} 页]    时间: {timestamp:<8} |    突变倍数: {current_ratio:.2f}x")
                        
                        diff_history = [] 

                if current_diff > 50:
                    diff_history.append(current_diff)
                    if len(diff_history) > 20:
                        diff_history.pop(0)
                
                pending_frame = frame.copy()
                last_gray = gray_blur
                
                progress = (count / total_frames) * 100
                bar = "█" * int(progress // 5) + "░" * (20 - int(progress // 5))
                print(f"\r  ⏳ 进度: [{bar}] {progress:4.1f}%   已截: {local_slide_count:02d} 页", end="")
            else:
                last_gray = gray_blur
                pending_frame = frame.copy()
            
        count += 1
    
    cap.release()
    print(f"\n✅ {video_path} 处理完毕，贡献了 {local_slide_count} 页。")
    return local_slide_count

if __name__ == "__main__":
    # 1. 扫描所有视频
    video_files = [f for f in os.listdir(".") if f.lower().endswith(('.mp4', '.mkv', '.avi'))]
    
    if not video_files:
        print("📁 没发现视频文件！请把 EXE 和视频放在同一个文件夹里。")
    else:
        # 2. 创建唯一的总 PPT 对象
        global_prs = Presentation()
        total_slides_all = 0
        
        print(f"🚀 发现 {len(video_files)} 个视频，准备合并到一个 PPT 中...")
        
        # 3. 循环处理每个视频，但往同一个 global_prs 里塞图片
        for idx, video in enumerate(video_files, 1):
            slides_added = process_one_video(video, global_prs, idx, len(video_files))
            total_slides_all += slides_added
        
        # 4. 所有视频跑完后，保存这个唯一的 PPT
        print(f"\n" + "="*65)
        print(f"💾 正在保存最终合并文件: {OUTPUT_FILENAME} ...")
        try:
            global_prs.save(OUTPUT_FILENAME)
            print(f"🎉 成功！所有视频共生成 {total_slides_all} 页，已保存在一个文件中。")
        except PermissionError:
            print(f"❌ 保存失败：请确保 {OUTPUT_FILENAME} 没有被打开！")

    # --- 防闪退 ---
    print("\n" + "·"*65)
    input("🍵 任务结束！按回车键（Enter）退出...")