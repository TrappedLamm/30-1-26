import cv2
import os
import numpy as np
from pptx import Presentation
from pptx.util import Inches

# ================= 配置区 =================
INTERVAL_SECONDS = 1      # 采样频率 (秒)
SENSITIVITY_FACTOR = 1.3  # 后期突变判定倍数
EARLY_THRESHOLD = 8000    # 前150秒硬阈值 (8000像素)
# ==========================================

def add_slide(prs, frame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    temp_img = "temp_capture.jpg"
    cv2.imwrite(temp_img, frame)
    slide.shapes.add_picture(temp_img, 0, 0, width=Inches(10))
    if os.path.exists(temp_img):
        os.remove(temp_img)

def video_to_ppt(video_path):
    output_name = os.path.splitext(video_path)[0] + "_精炼版.pptx"
    cap = cv2.VideoCapture(video_path)
    prs = Presentation()
    
    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    frame_interval = int(fps * INTERVAL_SECONDS)
    
    last_gray = None 
    pending_frame = None   
    diff_history = []      
    slide_count = 0
    count = 0

    print(f"\n" + "="*65)
    print(f"🎬 正在收割: {video_path}")
    print("="*65)

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret: 
            if pending_frame is not None:
                add_slide(prs, pending_frame)
            break
        
        if count % frame_interval == 0:
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            gray_blur = cv2.GaussianBlur(gray, (21, 21), 0)
            
            if last_gray is not None:
                delta = cv2.absdiff(last_gray, gray_blur)
                thresh = cv2.threshold(delta, 25, 255, cv2.THRESH_BINARY)[1]
                current_diff = cv2.countNonZero(thresh)
                
                # 计算突变倍数
                avg_diff = np.mean(diff_history) if len(diff_history) > 0 else current_diff
                current_ratio = current_diff / avg_diff if avg_diff > 0 else 0
                
                is_triggered = False

                # 判定逻辑
                if count < fps * 150:
                    if current_diff > EARLY_THRESHOLD:
                        is_triggered = True
                else:
                    if len(diff_history) > 3 and current_ratio > SENSITIVITY_FACTOR:
                        is_triggered = True

                if is_triggered:
                    if pending_frame is not None:
                        add_slide(prs, pending_frame)
                        slide_count += 1
                        # --- UI 优化：增加间距和分隔符，让数据更有呼吸感 ---
                        timestamp = f"{count//int(fps)}s"
                        print(f"\r✨ [第 {slide_count:02d} 页]    时间: {timestamp:<8} |    突变倍数: {current_ratio:.2f}x")
                        diff_history = [] 

                # 记录背景波动
                if current_diff > 50:
                    diff_history.append(current_diff)
                    if len(diff_history) > 20:
                        diff_history.pop(0)
                
                pending_frame = frame.copy()
                last_gray = gray_blur
                
                # --- 底部进度条：极致清爽 ---
                progress = (count / total_frames) * 100
                bar = "█" * int(progress // 5) + "░" * (20 - int(progress // 5))
                print(f"\r  ⏳ 进度: [{bar}] {progress:4.1f}%   已截: {slide_count:02d} 页", end="")
            else:
                last_gray = gray_blur
                pending_frame = frame.copy()
            
        count += 1
    
    try:
        prs.save(output_name)
        print(f"\n\n✅ 处理完毕，总计生成 {slide_count} 页课件。\n" + "="*65)
    except PermissionError:
        print(f"\n❌ 保存失败：请确保 PPT 文件未被打开！")
    cap.release()

if __name__ == "__main__":
    video_files = [f for f in os.listdir(".") if f.lower().endswith(('.mp4', '.mkv', '.avi'))]
    for v in video_files: 
        video_to_ppt(v)